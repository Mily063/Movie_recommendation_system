import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.xmlchemy import OxmlElement

def apply_slide_styling(slide):
    """Apply consistent styling to a slide."""
    # Apply light background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 250)  # Very light blue-gray

    # Style the title
    if slide.shapes.title:
        title_shape = slide.shapes.title
        title_frame = title_shape.text_frame
        title_frame.margin_bottom = Inches(0.08)
        title_frame.margin_top = Inches(0.08)

        for paragraph in title_frame.paragraphs:
            paragraph.font.size = Pt(32)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(44, 62, 80)  # Dark blue-gray
            paragraph.alignment = PP_ALIGN.CENTER

def apply_text_styling(text_frame, is_code=False):
    """Apply consistent styling to text."""
    text_frame.word_wrap = True
    text_frame.margin_bottom = Inches(0.05)
    text_frame.margin_top = Inches(0.05)
    text_frame.margin_left = Inches(0.1)
    text_frame.margin_right = Inches(0.1)

    for i, paragraph in enumerate(text_frame.paragraphs):
        if is_code:
            paragraph.font.name = 'Courier New'
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = RGBColor(41, 128, 185)  # Blue
        else:
            paragraph.font.name = 'Calibri'
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = RGBColor(52, 73, 94)  # Dark gray-blue

            # Add subtle styling for bullet points
            if paragraph.level >= 1:
                paragraph.font.size = Pt(16)
                if paragraph.level == 2:
                    paragraph.font.size = Pt(14)

def add_slide_transition(slide, transition_type='fade'):
    """
    Add a transition effect to a slide.
    Note: This is a placeholder function as python-pptx doesn't directly support transitions.
    In a real implementation, this would require XML manipulation.
    """
    # This is a placeholder - python-pptx doesn't directly support transitions
    # In a production environment, you would need to manipulate the XML directly
    pass

def create_title_slide(prs, title, subtitle=None):
    """Create a specially styled title slide."""
    slide_layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(slide_layout)

    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title

    # Set subtitle if provided
    if subtitle and hasattr(slide.placeholders, 'subtitle'):
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle

    # Apply special styling for title slide
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(235, 245, 255)  # Light blue background

    # Style the title
    title_frame = title_shape.text_frame
    title_frame.margin_bottom = Inches(0.1)
    title_frame.margin_top = Inches(0.1)

    for paragraph in title_frame.paragraphs:
        paragraph.font.size = Pt(44)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(41, 128, 185)  # Blue
        paragraph.alignment = PP_ALIGN.CENTER

    # Style the subtitle if it exists
    if subtitle and hasattr(slide.placeholders, 'subtitle'):
        subtitle_frame = subtitle_shape.text_frame
        for paragraph in subtitle_frame.paragraphs:
            paragraph.font.size = Pt(28)
            paragraph.font.italic = True
            paragraph.font.color.rgb = RGBColor(52, 73, 94)  # Dark gray-blue
            paragraph.alignment = PP_ALIGN.CENTER

    # Add a decorative line
    left = Inches(2)
    top = Inches(3.5)
    width = Inches(6)
    height = Inches(0.05)
    shape = slide.shapes.add_shape(
        1, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(41, 128, 185)  # Blue
    shape.line.fill.background()

    add_slide_transition(slide)

    return slide

def create_slide(prs, title, content):
    """Create a slide with title and content."""
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)

    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title

    # Set content
    content_shape = slide.shapes.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.clear()

    for line in content.split('\n'):
        p = text_frame.add_paragraph()
        if line.startswith('- '):
            p.text = line[2:]
            p.level = 1
        elif line.startswith('  - '):
            p.text = line[4:]
            p.level = 2
        else:
            p.text = line

    # Apply styling
    apply_slide_styling(slide)
    apply_text_styling(text_frame)
    add_slide_transition(slide)

    return slide

def create_code_slide(prs, title, code):
    """Create a slide with title and code."""
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)

    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title

    # Set code content
    content_shape = slide.shapes.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.clear()

    p = text_frame.add_paragraph()
    p.text = code

    # Apply styling
    apply_slide_styling(slide)
    apply_text_styling(text_frame, is_code=True)
    add_slide_transition(slide)

    return slide

def create_closing_slide(prs):
    """Create a closing 'Thank you' slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Apply background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(235, 245, 255)  # Light blue background

    # Add thank you text
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(1.5)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.add_paragraph()
    p.text = "Dziękuję za uwagę!"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(41, 128, 185)  # Blue
    p.alignment = PP_ALIGN.CENTER

    # Add a decorative line
    left = Inches(2)
    top = Inches(3.5)
    width = Inches(6)
    height = Inches(0.05)
    shape = slide.shapes.add_shape(
        1, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(41, 128, 185)  # Blue
    shape.line.fill.background()

    add_slide_transition(slide)

    return slide

def markdown_to_pptx(md_file, pptx_file):
    """Convert Markdown file to PowerPoint presentation."""
    prs = Presentation()

    # Set slide width and height (16:9 aspect ratio)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Read markdown file
    with open(md_file, 'r', encoding='utf-8') as f:
        md_content = f.read()

    # Split content by slide separator
    slides_content = md_content.split('---')

    # Process first slide as title slide if possible
    first_slide_processed = False

    for slide_idx, slide_content in enumerate(slides_content):
        if not slide_content.strip():
            continue

        # Extract title and content
        lines = slide_content.strip().split('\n')

        # Find the first heading
        title = ""
        subtitle = ""
        content_start = 0

        for i, line in enumerate(lines):
            if line.startswith('# '):
                title = line[2:]
                content_start = i + 1

                # Check if next line could be a subtitle
                if i + 1 < len(lines) and lines[i + 1].startswith('## '):
                    subtitle = lines[i + 1][3:]
                    content_start = i + 2

                break
            elif line.startswith('## '):
                title = line[3:]
                content_start = i + 1
                break
            elif line.startswith('### '):
                title = line[4:]
                content_start = i + 1
                break

        content = '\n'.join(lines[content_start:])

        # First slide becomes a title slide
        if slide_idx == 0 and not first_slide_processed:
            create_title_slide(prs, title, subtitle)
            first_slide_processed = True
            continue

        # Check if content contains code block
        if '```python' in content and '```' in content:
            # Extract code
            code_start = content.find('```python') + 10
            code_end = content.find('```', code_start)
            code = content[code_start:code_end].strip()

            # Extract content before code
            before_code = content[:content.find('```python')].strip()

            # Create slide with title and content before code
            if before_code:
                create_slide(prs, title, before_code)

            # Create slide with code
            create_code_slide(prs, title + " - Kod", code)
        else:
            # Create regular slide
            create_slide(prs, title, content)

    # Add a closing slide
    create_closing_slide(prs)

    # Save presentation
    prs.save(pptx_file)
    print(f"Presentation saved as {pptx_file}")

if __name__ == "__main__":
    md_file = "Prezentacja_Systemu_Rekomendacji_Filmow.md"
    pptx_file = "Prezentacja_Systemu_Rekomendacji_Filmow.pptx"

    if os.path.exists(md_file):
        markdown_to_pptx(md_file, pptx_file)
    else:
        print(f"File {md_file} not found.")
